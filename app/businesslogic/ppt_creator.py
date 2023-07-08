"""Module creating lyrics ppt"""
import os
import uuid

from pptx import Presentation

class LyricsPptCreator:
    """Class creating lyrics ppt file"""

    NEW_SECTION_TOKEN = '-'
    NUM_OF_LANGUAGES = 4
    EMPTY_LANGUAGE_TOKEN = '*'
    UNKNOWN_HYMN_NUM_TOKEN = '#'

    def __init__(self) -> None:
        pass

    def create_lyrics_ppt(self, input_txt_file_path) -> str:
        """Function creating lyrics ppt file"""

        ppt = Presentation(self.__get_template_file_path())
        lyrics_layout = ppt.slide_layouts[1]

        # Flip below boolean to true if you want to test with test_lyrics.txt file
        testing = False
        lyrics_file_path = ''
        if testing:
            lyrics_file_path = os.path.abspath('app/static/sample/lyrics/test_lyrics.txt')
        else:
            lyrics_file_path = input_txt_file_path
        self.__validate_lyrics_file(lyrics_file_path)

        hymn_nums = [0] * 5

        with open(lyrics_file_path, 'r', encoding='UTF-8') as lyrics_file:
            lines = lyrics_file.readlines()

            line_in_title_slide_counter = 0
            line_in_lyric_slide_counter = 0
            slide_counter = 0
            section_counter = 0
            for line in lines:
                # title slide
                if line_in_title_slide_counter < 4:
                    if len(line.strip()) == 0: # in case language is missing
                        hymn_nums[line_in_title_slide_counter] = self.EMPTY_LANGUAGE_TOKEN
                        line_in_title_slide_counter += 1
                        continue

                    hymn_num_name_list = line.split(' ', 1)
                    hymn_num = hymn_num_name_list[0].strip()
                    if hymn_num == '': # in case of " hymn_name"
                        hymn_nums[line_in_title_slide_counter] = self.UNKNOWN_HYMN_NUM_TOKEN
                    else: # in case of "# hymn_name" or "123 hymn_name"
                        hymn_nums[line_in_title_slide_counter] = hymn_num

                    ppt.slides[0].shapes[line_in_title_slide_counter].text = hymn_num
                    hymn_name = hymn_num_name_list[1].strip()
                    ppt.slides[0].shapes[line_in_title_slide_counter + 4].text = hymn_name
                    line_in_title_slide_counter += 1
                    continue

                # lyrics slide
                if self.__is_new_section(line):
                    if section_counter % 2 == 0: # this means we are on new slide
                        slide_counter += 1
                        ppt.slides.add_slide(lyrics_layout)
                        line_in_lyric_slide_counter = 0
                        section_counter = 0

                    line_split_list = line.split(' ', 1)
                    if len(line_split_list) == 2: # in case of "- 1"
                        ppt.slides[slide_counter].shapes[line_in_lyric_slide_counter].text = line_split_list[1].strip()
                    elif line_split_list[0][len(line_split_list)].isnumeric(): # in case of "-1"
                        ppt.slides[slide_counter].shapes[line_in_lyric_slide_counter].text = line_split_list[0][len(line_split_list)].strip()

                    line_in_lyric_slide_counter += 1
                    section_counter += 1
                    continue

                ppt.slides[slide_counter].shapes[line_in_lyric_slide_counter].text = line.strip()
                line_in_lyric_slide_counter += 1

        output_file_name = self.__get_output_file_name(hymn_nums)
        uuid_id = uuid.uuid4()
        os.mkdir(self.__get_output_file_directory(uuid_id), 0o777)
        output_file_path = self.__get_output_file_path(uuid_id, output_file_name)
        ppt.save(output_file_path)
        print(f'Successfully created a ppt file in path, {output_file_path}')

        return output_file_path

    def __get_template_file_path(self):
        working_dir = self.__get_current_directory()
        return os.path.join(working_dir, 'static/template/empty_template.pptx')

    def __get_output_file_directory(self, uuid_id):
        working_dir = self.__get_current_directory()
        return os.path.join(os.path.realpath(working_dir + '/../'), f'resource/output/{uuid_id}')

    def __get_output_file_path(self, uuid_id, output_file_name):
        return os.path.join(self.__get_output_file_directory(uuid_id), output_file_name)

    def __get_current_directory(self):
        return os.path.abspath(os.path.dirname(os.path.dirname(__file__)))

    def __validate_lyrics_file(self, lyrics_file_path):
        print(f'Validation of the lyrics file, \"{lyrics_file_path}\", is starting')

        with open(lyrics_file_path, 'r', encoding='UTF-8') as lyrics_file:
            lines = lyrics_file.readlines()
            line_in_section_counter = 0
            for line in lines:
                if self.__is_new_section(line):
                    if line_in_section_counter == self.NUM_OF_LANGUAGES:
                        line_in_section_counter = 0
                        continue
                    raise ValueError(
                        f'The number of lines between \"{self.NEW_SECTION_TOKEN}\" \
should always be {self.NUM_OF_LANGUAGES}, but it is {line_in_section_counter} in a section')

                line_in_section_counter += 1

        print(f'Validation of the lyrics file, \"{lyrics_file_path}\", is completed successfully')

    def __is_new_section(self, line) -> bool:
        return line.startswith(self.NEW_SECTION_TOKEN)

    def __get_output_file_name(self, hymn_nums) -> str:
        output_file_name = ''
        if hymn_nums[2] != '*':
            output_file_name += f'K{hymn_nums[2]}'

        if hymn_nums[0] != '*':
            output_file_name += f'_E{hymn_nums[0]}'

        if hymn_nums[1] != '*':
            output_file_name += f'_S{hymn_nums[1]}'

        if hymn_nums[3] != '*':
            output_file_name += f'_C{hymn_nums[3]}'

        return output_file_name + '.pptx'
