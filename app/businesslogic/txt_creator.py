"""Module creating lyrics txt file from a pptx file"""
import os
import uuid

from pathlib import Path
from pptx import Presentation

class LyricsTxtCreator:
    """Class creating lyrics txt file"""

    def __init__(self) -> None:
        pass

    def create_lyrics_txt(self, input_pptx_file_path, uuid_id) -> str:
        """Function creating lyrics txt file"""

        ppt = Presentation(input_pptx_file_path)

        output_file_name = f'{Path(input_pptx_file_path).stem}.txt'
        output_directory = self.__get_output_file_directory(uuid_id)
        if not os.path.exists(output_directory):
            os.mkdir(self.__get_output_file_directory(uuid_id), 0o777)
        output_file_path = self.__get_output_file_path(uuid_id, output_file_name)

        with open(output_file_path, 'a') as txt_file:
            for slide in ppt.slides:
                for shape in slide.shapes:
                    txt_file.write(f'{shape.text}\n')

        return output_file_path

    def __get_output_file_directory(self, uuid_id):
        working_dir = self.__get_current_directory()
        return os.path.join(os.path.realpath(working_dir + '/../'), f'resource/output/{uuid_id}')

    def __get_output_file_path(self, uuid_id, output_file_name):
        return os.path.join(self.__get_output_file_directory(uuid_id), output_file_name)

    def __get_current_directory(self):
        return os.path.abspath(os.path.dirname(os.path.dirname(__file__)))