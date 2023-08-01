"""Module creating lyrics ppt"""
import os

from pathlib import Path
from pptx import Presentation

class WedSermonLyricsPptCreator:
    """Class creating lyrics ppt file for wednesday sermon"""

    NEW_SECTION_TOKEN = '-'
    NUM_OF_LANGUAGES = 4
    EMPTY_LANGUAGE_TOKEN = '*'
    UNKNOWN_HYMN_NUM_TOKEN = '#'

    def __init__(self) -> None:
        pass

    def create_lyrics_ppt(self, input_txt_file_path, uuid_id) -> str:
        """Function creating lyrics ppt file"""

        ppt = Presentation(self.__get_template_file_path())
        main_layout = ppt.slide_layouts[0]

        with open(input_txt_file_path, 'r', encoding='UTF-8') as lyrics_file:
            lines = lyrics_file.readlines()

            slide_counter = 0
            for line in lines:
                text_shape = ppt.slides[slide_counter].shapes[0]
                text_shape.text = line.strip()
                text_frame = text_shape.text_frame
                working_dir = self.__get_current_directory()
                font_path = os.path.join(os.path.realpath(working_dir + '/../'), 'app/static/fonts/arial.ttf')
                text_frame.fit_text(font_family=None, max_size=30, bold=None, italic=None, font_file=font_path)
                text_frame.auto_size = None
                slide_counter += 1
                ppt.slides.add_slide(main_layout)

        output_file_name = f'{Path(input_txt_file_path).stem}.pptx'
        output_directory = self.__get_output_file_directory(uuid_id)
        if not os.path.exists(output_directory):
            os.mkdir(self.__get_output_file_directory(uuid_id), 0o777)
        output_file_path = self.__get_output_file_path(uuid_id, output_file_name)

        ppt.save(output_file_path)
        print(f'Successfully created a ppt file in path, {output_file_path}')

        return output_file_path

    def __get_template_file_path(self):
        working_dir = self.__get_current_directory()
        return os.path.join(working_dir, 'static/template/wed_sermon_template.pptx')

    def __get_output_file_directory(self, uuid_id):
        working_dir = self.__get_current_directory()
        return os.path.join(os.path.realpath(working_dir + '/../'), f'resource/output/{uuid_id}')

    def __get_output_file_path(self, uuid_id, output_file_name):
        return os.path.join(self.__get_output_file_directory(uuid_id), output_file_name)

    def __get_current_directory(self):
        return os.path.abspath(os.path.dirname(os.path.dirname(__file__)))
