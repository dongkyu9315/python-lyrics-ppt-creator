"""Module creating lyrics ppt"""
import os
import uuid

from pptx import Presentation

class LyricsPptCreator:
    """Class creating lyrics ppt file"""

    NEW_SECTION_TOKEN = '-'
    NUM_OF_LANGUAGES = 4

    def __init__(self) -> None:
        pass

    def create_lyrics_ppt(self, input_file_path) -> str:
        """Function creating lyrics ppt file"""

        ppt = Presentation(self.__get_template_file_path())
        lyrics_layout = ppt.slide_layouts[1]

        # Flip below boolean to true if you want to test with test_lyrics.txt file
        testing = False
        lyrics_file_path = ''
        if testing:
            lyrics_file_path = os.path.abspath('resource/lyrics/test_lyrics.txt')
        else:
            lyrics_file_path = input_file_path
        self.__validate_lyrics_file(lyrics_file_path)

        hymn_nums = [0] * 5

        with open(lyrics_file_path, 'r', encoding='UTF-8') as lyrics_file:
            lines = lyrics_file.readlines()

            line_in_title_slide_counter = 0
            line_in_lyric_slide_counter = 0
            title_slide_counter = 0
            slide_counter = 0
            section_counter = 0
            for line in lines:
                # title slide
                if title_slide_counter < 4:
                    hymn_num = line.split(' ', 1)[0].strip()
                    hymn_nums[line_in_title_slide_counter] = hymn_num
                    ppt.slides[0].shapes[line_in_title_slide_counter].text = hymn_num
                    hymn_name = line.split(' ', 1)[1].strip()
                    ppt.slides[0].shapes[line_in_title_slide_counter + 4].text = hymn_name
                    line_in_title_slide_counter += 1
                    title_slide_counter += 1
                    continue

                # lyrics slide
                if self.__is_new_section(line):
                    if section_counter % 2 == 0:
                        slide_counter += 1
                        ppt.slides.add_slide(lyrics_layout)
                        line_in_lyric_slide_counter = 0
                        section_counter = 0

                        line_split_list = line.split(' ', 1)
                        if len(line_split_list) == 2:
                            ppt.slides[slide_counter].shapes[8].text = line_split_list[1].strip()
                            line_in_lyric_slide_counter = 0

                    section_counter += 1
                    continue

                ppt.slides[slide_counter].shapes[line_in_lyric_slide_counter].text = line.strip()
                line_in_lyric_slide_counter += 1

        output_file_name = f'E{hymn_nums[0]}_K{hymn_nums[1]}_S{hymn_nums[2]}_C{hymn_nums[3]}.pptx'
        uuid_id = uuid.uuid4()
        os.mkdir(self.__get_output_file_directory(uuid_id), 0o777)
        output_file_path = self.__get_output_file_path(uuid_id, output_file_name)
        ppt.save(output_file_path)
        print(f'Successfully created a ppt file in path, {output_file_path}')

        return output_file_path

    def __get_template_file_path(self):
        return os.path.abspath('resource/template/empty_template.pptx')

    def __get_output_file_directory(self, uuid_id):
        return os.path.abspath(f'resource/output/{uuid_id}')

    def __get_output_file_path(self, uuid_id, output_file_name):
        return os.path.join(self.__get_output_file_directory(uuid_id), output_file_name)

    def __validate_lyrics_file(self, lyrics_file_path):
        print(f'Validation of the lyrics file, \"{lyrics_file_path}\", is starting')

        with open(lyrics_file_path, 'r', encoding='UTF-8') as lyrics_file:
            lines = lyrics_file.readlines()
            line_in_section_counter = 0
            for line in lines:
                if self.__is_new_section(line):
                    if line_in_section_counter == self.NUM_OF_LANGUAGES:
                        line_in_section_counter = 0
                        continue
                    raise ValueError(
                        f'The number of lines between \"{self.NEW_SECTION_TOKEN}\" \
should always be {self.NUM_OF_LANGUAGES}, but it is {line_in_section_counter} in a section')

                line_in_section_counter += 1

        print(f'Validation of the lyrics file, \"{lyrics_file_path}\", is completed successfully')

    def __is_new_section(self, line) -> bool:
        return line.startswith(self.NEW_SECTION_TOKEN)
